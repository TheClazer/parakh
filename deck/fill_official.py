#!/usr/bin/env python3
"""Fill the official Redrob Track-1 submission template with Parakh, in-place."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x20, 0x21, 0x24)
PUR = RGBColor(0x6D, 0x28, 0xD9)
SLA = RGBColor(0x5F, 0x63, 0x68)
LAV = RGBColor(0xF5, 0xF1, 0xFE)
TEAL = RGBColor(0x0B, 0x7A, 0x60)
CLAY = RGBColor(0xB3, 0x3A, 0x22)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
F = "Arial"

prs = Presentation("Parakh_official.pptx")
S = list(prs.slides)


def shp(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_qa(slide, sid, pairs, qsz=11, asz=10.5, gap=6):
    tf = shp(slide, sid).text_frame
    tf.clear(); tf.word_wrap = True
    first = True
    for q, a in pairs:
        pq = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        pq.space_after = Pt(1)
        r = pq.add_run(); r.text = q
        r.font.name = F; r.font.size = Pt(qsz); r.font.bold = True; r.font.color.rgb = PUR
        pa = tf.add_paragraph(); pa.space_after = Pt(gap); pa.line_spacing = 1.03
        r = pa.add_run(); r.text = a
        r.font.name = F; r.font.size = Pt(asz); r.font.color.rgb = INK


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(para.get("after", 0))
        p.line_spacing = para.get("line", 1.0)
        for t, sz, c, b in para["runs"]:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b
    return tb


def box(slide, x, y, w, h, label, sz=9, fill=LAV, line=PUR):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = line; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    try: b.adjustments[0] = 0.14
    except Exception: pass
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 0.95
    r = p.add_run(); r.text = label
    r.font.name = F; r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = INK
    return b


def arrow(slide, x, y, w=0.28, h=0.6):
    add_text(slide, x, y, w, h, [{"runs": [("→", 13, PUR, True)]}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def chain(slide, y, items, x0=0.42, x1=9.7, h=0.66, sz=8.5):
    n = len(items); aw = 0.24
    bw = (x1 - x0 - aw * (n - 1)) / n
    x = x0
    for i, it in enumerate(items):
        box(slide, x, y, bw, h, it, sz=sz)
        x += bw
        if i < n - 1:
            arrow(slide, x - 0.02, y, aw + 0.04, h)
            x += aw


def fill_field(slide, sid, value, sz=15):
    tf = shp(slide, sid).text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = " " + value
    r.font.name = F; r.font.size = Pt(sz); r.font.bold = True; r.font.color.rgb = PUR


# ================= S1 - title fields =================
fill_field(S[0], 55, "Evolve")
fill_field(S[0], 57, "Rayyan Ahmed Shaikh")
fill_field(S[0], 56, "Track 1 - The Data & AI Challenge: Intelligent Candidate Discovery & Ranking", sz=13)
add_text(S[0], 0.34, 4.85, 9.3, 0.5,
         [{"runs": [("PARAKH", 15, PUR, True),
                    ("  (परख - to assay; to test whether gold is genuine)  ·  Rayyan Ahmed Shaikh · Inchara · Harshita Nagesh", 12, INK, False)]}])

# ================= S2 - Solution Overview =================
set_qa(S[1], 64, [
    ("What is your proposed solution?",
     "Parakh - a ranker that verifies every profile before trusting a word of it. An Integrity Gate rejects the dataset's planted fakes outright; two frontier LLMs (DeepSeek-V4 + Qwen3-235B) grade every plausible candidate against the JD offline; a third model family (GLM-5.2) re-judges the top; only candidates all three families endorse reach the top-10. Every model call is pre-computed and cached - the submitted ranker runs on a laptop CPU in ~3 minutes, fully offline, and reproduces the CSV byte-for-byte."),
    ("What differentiates it from traditional candidate matching?",
     "Traditional matchers trust what candidates claim and count keyword overlap - and this dataset is engineered to kill exactly that (stuffed skill lists, ~80 planted honeypots, real gems hidden in plain language). Parakh ranks on verified evidence instead: what the career narrative proves, what 23 behavioral signals confirm, and what three independent AI judges agree on. Assay, don't match."),
], asz=11)

# ================= S3 - JD Understanding =================
set_qa(S[2], 71, [
    ("Key requirements extracted from the JD",
     "Must-haves: production embedding retrieval (sentence-transformers / BGE / E5), vector-DB or hybrid-search operations (FAISS, Pinecone, Elasticsearch), rigorous ranking evaluation (NDCG, MRR, A/B), strong Python, 5-9 yrs (ideal 6-8) mostly at product companies, India-based or willing to relocate. Explicit disqualifiers the JD itself states: consulting-only careers, research without production, recent LangChain-only 'AI experience', CV/speech-only depth, title-chasers, unreachable candidates."),
    ("Which candidate signals matter most?",
     "The skills array is the least trustworthy field - it is what stuffers stuff. We weight the career narrative (did they build and ship retrieval / ranking / recsys), title trajectory, product-vs-services history, and behavioral signals: recruiter response rate, last-active recency, notice period, open-to-work. A perfect profile that never replies is not a hire."),
    ("Evaluating fit beyond keyword matching",
     "Free-text evidence extraction plus LLM judgment: 7,334 candidates describe real ranking work under non-ML titles, and keyword filters miss every one of them. Our judges read the story, not the tags - and a recall-rescue sweep re-checks all 100k for narrative evidence the rules under-ranked, which recovered 9 hidden gems."),
], asz=10)

# ================= S4 - Ranking Methodology =================
set_qa(S[3], 78, [
    ("Retrieve, score, rank",
     "Stage 1 - Integrity Gate: reject impossible profiles ('expert' skills used 0 months, tenure longer than the whole career) and keyword-stuffers (AI skills with zero supporting narrative). Stage 2 - a transparent JD rubric scores all 100,000: role, narrative evidence, must-have tech, experience band, company type, location, availability. Stage 3 - offline judging: DeepSeek-V4 and Qwen3-235B each grade the rule-shortlist top-3,000 PLUS a 2,500-profile full-pool sweep; GLM-5.2 re-judges the top-100."),
    ("Models, algorithms, heuristics",
     "Frontier open-weight LLMs as offline judges (via Nebius Token Factory), deterministic feature engineering, and narrative-evidence extraction. No fine-tuning required. The submitted ranker is pure Python reading cached judgments: fast, auditable, reproducible."),
    ("Combining signals into one ranking",
     "final score = average judge tier (0-5) + rule score as the within-tier tiebreak, scaled to [0,1]. Honeypots are forced to the floor; equal scores tie-break by candidate_id ascending, exactly as the spec requires. Judgment dominates; rules order within agreement bands."),
], asz=10)

# ================= S5 - Explainability =================
set_qa(S[4], 85, [
    ("How are ranking decisions explained?",
     "Every one of the 100 rows carries a 1-2 line justification citing real profile facts (title, years, named tech, response rate) tied to a JD requirement - with tone matched to rank: low ranks lead with the honest concern, not praise."),
    ("Preventing hallucinated justifications",
     "Reasons are generated strictly from the profile JSON, then programmatically verified: every tech or company token quoted must literally exist in that profile; failures are regenerated or replaced with grounded fallbacks. 95 of 100 final reasons pass strict verification; the rest use verified teacher text."),
    ("Handling inconsistent or suspicious profiles",
     "That is the Integrity Gate's job: 42 hard honeypots detected analytically, zero in our top-100 (the disqualification line is 10%). Stuffers are suppressed multiplicatively. We also tested a tempting extra check (last-active before signup) and rejected it - it fires on 7,496 profiles of generator noise. Precision over paranoia."),
], asz=10)

# ================= S6 - End-to-End Workflow =================
set_qa(S[5], 92, [
    ("From JD input to ranked output - one command",
     "python rank.py --candidates candidates.jsonl --out submission.csv   ·   ~3 min on a laptop CPU, no network, no GPU. The heavy thinking happened once, offline; artifacts ship with the repo."),
], asz=10.5)
chain(S[5], 3.15, ["candidates.jsonl\n(100,000 profiles)", "Integrity Gate\nfakes -> floor", "JD rubric scores\nall 100,000", "Cached 3-judge\ntiers merge"], h=0.78)
chain(S[5], 4.35, ["Blend: judge avg\n+ rule tiebreak", "Top-100 with\nverified reasons", "submission.csv\n+ XLSX (validated)"], h=0.78)

# ================= S7 - System Architecture =================
add_text(S[6], 0.53, 1.42, 9.0, 0.4,
         [{"runs": [("Two planes: the intelligence is pre-computed; the ranking step just reads it.", 12, SLA, False)]}])
box(S[6], 0.42, 1.95, 2.05, 1.62, "PRE-COMPUTE\n(offline, network OK)\nruns once, ships\nas artifacts", sz=9.5, fill=LAV)
chain(S[6], 2.05, ["Full-pool sweep\n+ features"], x0=2.75, x1=4.85, h=0.62)
chain(S[6], 2.05, ["DeepSeek-V4 + Qwen-235B\nteachers grade 5,500 each"], x0=5.05, x1=7.6, h=0.62)
chain(S[6], 2.05, ["GLM-5.2\n3rd judge, top-100"], x0=7.8, x1=9.7, h=0.62)
chain(S[6], 2.9, ["verified reasons\n(fact-checked)"], x0=2.75, x1=4.85, h=0.62)
chain(S[6], 2.9, ["four .jsonl artifacts, committed to the repo"], x0=5.05, x1=9.7, h=0.62)
box(S[6], 0.42, 4.75, 2.05, 0.62, "RANKING STEP\nCPU · no network", sz=9.5, fill=LAV, line=TEAL)
chain(S[6], 4.75, ["load artifacts", "gate + score\n100k", "blend +\ntiebreak", "top-100 CSV\n+ reasons"], x0=2.75, x1=9.7, h=0.62)
add_text(S[6], 0.53, 3.85, 9.0, 0.6,
         [{"runs": [("Self-test on every run: ", 10.5, INK, True),
                    ("~180-265 s wall (limit 300) · 0.21 GB RAM (limit 16) · bare command reproduces the submission byte-identically.", 10.5, SLA, False)]}])

# ================= S8 - Results & Performance =================
set_qa(S[7], 105, [
    ("Ranking quality",
     "Judged by an independent model that never touched the ranking (Llama-3.3-70B): NDCG@10 rose 0.64 (rules) -> 0.88 (+teacher) -> 1.00 (3-family ensemble); the final top-10 is unanimously tier-5 across DeepSeek, Qwen and GLM. Zero honeypots and zero stuffers in the top-100. The full-pool sweep rescued 9 genuine hidden gems the shortlist missed. Honest note: judge agreement is a strong proxy, not the organisers' hidden truth."),
    ("Runtime and compute constraints",
     "Built-in self-test prints proof on every run: ~180-265 s wall against the 300 s limit, 0.21 GB RAM against 16 GB, CPU-only, zero network calls. The bare reproduce command regenerates the submitted CSV byte-for-byte from the repo."),
], asz=10.5)
for i, (k, v, c) in enumerate([("1.00", "NDCG@10 (independent judge)", TEAL),
                               ("10/10", "top-10 tier-5, 3 model families", PUR),
                               ("0", "honeypots in our top-100", CLAY)]):
    x = 0.53 + i * 3.1
    b = box(S[7], x, 4.15, 2.9, 0.95, "", fill=LAV)
    add_text(S[7], x + 0.15, 4.25, 2.6, 0.75,
             [{"runs": [(k, 20, c, True)], "after": 1},
              {"runs": [(v, 9, SLA, False)]}])

# ================= S9 - Technologies =================
set_qa(S[8], 112, [
    ("Stack and why",
     "Ranker: Python 3.11, stdlib-first (orjson fast parse, openpyxl for XLSX) - deliberately dependency-light so Stage-3 reproduction cannot break. Offline judging: Nebius Token Factory serving DeepSeek-V4-Pro and Qwen3-235B as teachers, GLM-5.2 as third judge, Llama-3.3-70B as the independent evaluator - four different model families, chosen so that agreement actually means something. Sandbox: Gradio on Hugging Face Spaces, running the exact same code path. Git with full iteration history (12 commits of real evolution, not a final dump)."),
], asz=10.5)

# ================= S10 - Submission Assets =================
set_qa(S[9], 119, [
    ("Everything needed to verify us",
     ""),
], asz=10.5)
assets = [
    ("GitHub (full code + artifacts)", "github.com/TheClazer/parakh"),
    ("Live sandbox (Hugging Face Space)", "huggingface.co/spaces/TheClazer/parakh-ranker"),
    ("Ranked output", "Parakh_top100.xlsx / submission.csv - passes the official validator"),
    ("Reproduce", "python rank.py --candidates candidates.jsonl --out submission.csv"),
    ("Contact", "Rayyan Ahmed Shaikh · therayyn16@gmail.com · Team Evolve"),
]
for i, (k, v) in enumerate(assets):
    y = 2.0 + i * 0.62
    add_text(S[9], 0.75, y, 3.1, 0.5, [{"runs": [(k, 11.5, PUR, True)]}])
    add_text(S[9], 3.95, y, 5.6, 0.5, [{"runs": [(v, 11.5, INK, False)]}])

# ================= S11 - closing =================
add_text(S[10], 0.5, 4.85, 9.0, 0.5,
         [{"runs": [("PARAKH - assay, don't match  ·  Team Evolve  ·  thank you", 13, WHT, True)]}],
         align=PP_ALIGN.CENTER)

prs.save("Parakh_official.pptx")
print("filled", len(prs.slides._sldIdLst), "slides")
