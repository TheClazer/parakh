#!/usr/bin/env python3
"""Builds parakh_deck.pptx — Parakh approach deck for India Runs Track 1 (Team Evolve).

Design language: an assayer's hallmark. Deep-ink title/close slides, light content
slides, brass as the single sharp accent, teal = genuine, clay = fake.
Motif: a small brass "hallmark punch" ring, repeated as section markers.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ---------------------------------------------------------------
INK      = RGBColor(0x1B, 0x2A, 0x41)
INK_SOFT = RGBColor(0x33, 0x45, 0x5F)
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)
BRASS    = RGBColor(0xB0, 0x89, 0x4A)
BRASS_D  = RGBColor(0x8A, 0x69, 0x33)
TEAL     = RGBColor(0x0E, 0x7C, 0x66)
CLAY     = RGBColor(0xB4, 0x47, 0x2E)
SLATE    = RGBColor(0x5A, 0x64, 0x72)
PANEL    = RGBColor(0xF4, 0xF5, 0xF7)
TEAL_LT  = RGBColor(0xE4, 0xF0, 0xEB)
CLAY_LT  = RGBColor(0xF7, 0xE9, 0xE3)
BRASS_LT = RGBColor(0xF3, 0xEC, 0xDC)
MIST     = RGBColor(0xCF, 0xD6, 0xDE)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"
CREDITS = "Rayyan Ahmed Shaikh   ·   Inchara   ·   Harshita Nagesh"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def shape(s, kind, x, y, w, h, fill=None, line=None, lw=1.0, radius=None):
    shp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp


def card(s, x, y, w, h, fill=PANEL, radius=0.06):
    return shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, radius=radius)


def text(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "before" in para: p.space_before = Pt(para["before"])
        p.space_after = Pt(para.get("after", 0))
        if "line" in para: p.line_spacing = para["line"]
        for run in para["runs"]:
            r = p.add_run(); r.text = run[0]
            f = r.font
            f.size = Pt(run[1]); f.color.rgb = run[2]
            f.bold = run[3] if len(run) > 3 else False
            f.italic = run[4] if len(run) > 4 else False
            f.name = run[5] if len(run) > 5 else BODY
    return tb


def hallmark(s, cx, cy, d=0.30, color=BRASS):
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    ring.fill.background(); ring.line.color.rgb = color; ring.line.width = Pt(1.5)
    ring.shadow.inherit = False
    dd = d * 0.34
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - dd / 2), Inches(cy - dd / 2), Inches(dd), Inches(dd))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    dot.shadow.inherit = False


def kicker(s, x, y, label, color=BRASS):
    hallmark(s, x + 0.13, y + 0.12, 0.26, color)
    text(s, x + 0.42, y - 0.03, 9.0, 0.4, [{"runs": [(label.upper(), 13, color, True)]}])


def title(s, y, t, size=33):
    text(s, 0.7, y, 12.0, 1.0, [{"runs": [(t, size, INK, True, False, HEAD)]}])


def arrow(s, x, y, w=0.5, h=0.88, color=BRASS):
    text(s, x, y, w, h, [{"align": PP_ALIGN.CENTER, "runs": [("→", 20, color, True)]}], anchor=MSO_ANCHOR.MIDDLE)


def minibox(s, x, y, w, h, label, line_color, fill=PAPER):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line_color, lw=1.25, radius=0.1)
    text(s, x + 0.12, y, w - 0.24, h, [{"align": PP_ALIGN.CENTER, "line": 1.0, "runs": [(label, 12.5, INK, True)]}], anchor=MSO_ANCHOR.MIDDLE)


# ============================================================ 1 — TITLE
s = slide(INK)
for (cx, cy, d) in [(11.7, 1.5, 2.9), (12.7, 4.6, 1.7), (10.9, 5.9, 1.1)]:
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    ring.fill.background(); ring.line.color.rgb = INK_SOFT; ring.line.width = Pt(1.5); ring.shadow.inherit = False
hallmark(s, 0.95, 1.35, 0.4, BRASS)
text(s, 1.35, 1.15, 8, 0.5, [{"runs": [("THE DATA & AI CHALLENGE  ·  REDROB × INDIA RUNS", 13, MIST, True)]}])
text(s, 0.9, 2.05, 9.5, 1.5, [{"runs": [("Parakh", 76, PAPER, True, False, HEAD)]}])
text(s, 0.95, 3.42, 9.2, 0.6, [{"runs": [("/ pə-rəkh /  —  to assay; to test whether gold is genuine.", 20, BRASS, False, True, HEAD)]}])
text(s, 0.95, 4.35, 9.0, 1.0,
     [{"runs": [("A candidate ranker that ", 22, MIST), ("verifies", 22, PAPER, True, True),
                (" each profile before it trusts it —", 22, MIST)]},
      {"runs": [("telling genuine engineers from convincing fakes.", 22, MIST)]}])
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(6.0), Inches(3.2), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = BRASS; ln.line.fill.background(); ln.shadow.inherit = False
text(s, 0.95, 6.2, 11, 0.7,
     [{"runs": [("Team Evolve", 16, PAPER, True)]},
      {"before": 2, "runs": [(CREDITS, 13.5, MIST)]}])

# ============================================================ 2 — PROBLEM
s = slide(PAPER)
kicker(s, 0.7, 0.6, "The problem")
text(s, 0.7, 1.05, 7.4, 1.4, [{"runs": [("Great candidates are missed —", 34, INK, True, False, HEAD)]},
                              {"runs": [("not absent, just unseen.", 34, INK, True, False, HEAD)]}])
text(s, 0.7, 2.9, 6.7, 2.6,
     [{"line": 1.15, "after": 10, "runs": [("Recruiters scan hundreds of profiles and still miss the right person — because ", 16, SLATE), ("keyword filters can't see fit.", 16, INK, True)]},
      {"line": 1.15, "after": 10, "runs": [("They reward whoever lists the most buzzwords, and overlook the engineer who ", 16, SLATE), ("built the exact system", 16, INK, True), (" but described it plainly.", 16, SLATE)]},
      {"line": 1.15, "runs": [("In a pool of 100,000 the real fits are a needle in a haystack — and half the score rides on the top ten.", 16, SLATE)]}])
sx = 8.15
for i, (num, lab, col) in enumerate([("100,000", "profiles in the pool", INK),
                                     ("~ few hundred", "genuinely fit the role", BRASS_D),
                                     ("Top 10", "carry 50% of the score", TEAL)]):
    yy = 1.15 + i * 1.72
    card(s, sx, yy, 4.45, 1.5, PANEL)
    text(s, sx + 0.35, yy + 0.18, 3.9, 0.85, [{"runs": [(num, 38, col, True, False, HEAD)]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, sx + 0.37, yy + 1.0, 3.9, 0.4, [{"runs": [(lab, 14, SLATE)]}])

# ============================================================ 3 — TRAPS
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Why naïve matching loses")
title(s, 1.05, "The dataset is engineered to punish keyword search")
traps = [("Keyword stuffers", CLAY, CLAY_LT, "Lists every AI skill; the actual job is Marketing or Ops. Counting keywords ranks the fake at the top."),
         ("Plain-language gems", TEAL, TEAL_LT, "7,334 people describe real retrieval / ranking work under a non-ML title. Title filters miss them entirely."),
         ("Behavioural twins", BRASS_D, BRASS_LT, "Near-identical on paper; they differ only in engagement signals — who actually responds and shows up."),
         ("Honeypots (~80)", CLAY, CLAY_LT, "Impossible profiles (expert skills, 0 months' use). Rank >10% of them in the top 100 and you're disqualified.")]
gw, gh = 5.83, 2.15
for i, (t, ac, tint, desc) in enumerate(traps):
    x = 0.7 + (i % 2) * (gw + 0.24); y = 2.1 + (i // 2) * (gh + 0.24)
    card(s, x, y, gw, gh, tint)
    hallmark(s, x + 0.5, y + 0.55, 0.34, ac)
    text(s, x + 0.95, y + 0.32, gw - 1.2, 0.5, [{"runs": [(t, 20, INK, True, False, HEAD)]}])
    text(s, x + 0.5, y + 1.05, gw - 0.9, 1.0, [{"line": 1.12, "runs": [(desc, 14.5, SLATE)]}])

# ============================================================ 4 — WHAT EVERYONE ELSE BUILDS
s = slide(PAPER)
kicker(s, 0.7, 0.6, "The trap everyone falls into", CLAY)
title(s, 1.05, "What most teams will submit — and why it loses")
card(s, 0.7, 2.1, 5.83, 3.35, PANEL)
text(s, 1.05, 2.35, 5.2, 0.5, [{"runs": [("The obvious pipeline", 19, INK, True, False, HEAD)]}])
for i, b in enumerate(["Embed the job description", "Embed all 100,000 profiles", "Cosine similarity → sort"]):
    yy = 3.0 + i * 0.72
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.05, yy, 5.1, 0.58, fill=PAPER, line=MIST, lw=1.25, radius=0.14)
    text(s, 1.25, yy, 4.8, 0.58, [{"runs": [(f"{i+1}.  ", 13, SLATE, True), (b, 13.5, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
card(s, 6.8, 2.1, 5.83, 3.35, CLAY_LT)
text(s, 7.15, 2.35, 5.2, 0.5, [{"runs": [("What the traps do to it", 19, INK, True, False, HEAD)]}])
for i, b in enumerate(["Keyword-stuffers rank #1 — buzzwords beat real work",
                       "Plain-language gems get buried below them",
                       "Impossible honeypots slip into the top 10",
                       "→ honeypot rate breaks 10% → disqualified"]):
    hl = i == 3
    text(s, 7.15, 3.0 + i * 0.55, 5.2, 0.5,
         [{"runs": [("·  ", 15, CLAY, True), (b, 14, CLAY if hl else INK_SOFT, hl)]}])
text(s, 0.7, 5.75, 11.9, 0.5, [{"runs": [("The dataset was built to beat exactly this. So we don't play that game.", 15, INK, True, True)]}])

# ============================================================ 5 — THESIS
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Our approach")
title(s, 1.05, "Assay, don't match.", 40)
text(s, 0.7, 2.0, 11.9, 0.9,
     [{"line": 1.2, "runs": [("Like an assayer testing gold, Parakh ", 18, SLATE), ("verifies every profile before it trusts a word of it", 18, INK, True),
                             (" — then ranks only on evidence that survives.", 18, SLATE)]}])
principles = [("Read the work, not the words", "Evidence in the career story outweighs any skills list."),
              ("Reject the impossible", "Fakes and honeypots are removed before ranking, not after."),
              ("Trust what judges agree on", "The very top is where independent AI graders converge."),
              ("Prove it, don't hope", "Every design choice is measured against an independent judge.")]
for i, (t, d) in enumerate(principles):
    x = 0.7 + (i % 2) * 6.0; y = 3.15 + (i // 2) * 1.55
    hallmark(s, x + 0.28, y + 0.26, 0.36, BRASS)
    text(s, x + 0.78, y, 5.0, 0.5, [{"runs": [(t, 18, INK, True, False, HEAD)]}])
    text(s, x + 0.78, y + 0.5, 5.05, 0.9, [{"line": 1.1, "runs": [(d, 14, SLATE)]}])

# ============================================================ 6 — THE TRICK
s = slide(PAPER)
kicker(s, 0.7, 0.6, "The core idea")
title(s, 1.05, "An LLM's judgment — at CPU speed")
card(s, 0.7, 2.1, 11.93, 1.15, INK)
text(s, 1.05, 2.1, 11.2, 1.15,
     [{"runs": [("The hard rule:  ", 15, BRASS, True), ("the ranking step must run offline, CPU-only, no LLM calls, under 5 minutes over 100,000 profiles. "
                 "An LLM call per candidate simply cannot fit — or scale in production.", 15, PAPER)]}],
     anchor=MSO_ANCHOR.MIDDLE)
steps = [("Think — offline", "A frontier LLM (DeepSeek-V4) grades the plausible candidates on the job rubric. Slow, brilliant, unconstrained.", BRASS_D),
         ("Bake — into artifacts", "Its verdicts are frozen into a small cached file that ships with the code — a precomputed feature.", INK_SOFT),
         ("Serve — in milliseconds", "The offline ranker reads those verdicts. No LLM, no GPU, no network at run-time.", TEAL)]
for i, (t, d, c) in enumerate(steps):
    x = 0.7 + i * 4.05
    card(s, x, 3.55, 3.78, 2.0, PANEL)
    text(s, x + 0.32, 3.75, 3.2, 0.5, [{"runs": [(t, 17, c, True, False, HEAD)]}])
    text(s, x + 0.32, 4.28, 3.2, 1.2, [{"line": 1.12, "runs": [(d, 13, SLATE)]}])
    if i < 2:
        arrow(s, x + 3.82, 3.55, 0.28, 2.0, BRASS)
text(s, 0.7, 5.8, 11.9, 0.5, [{"runs": [("The constraint becomes the advantage: ", 15, INK, True), ("frontier intelligence, distilled into a ranker that's fast, legal, and reproducible.", 15, SLATE)]}])

# ============================================================ 7 — PIPELINE
s = slide(PAPER)
kicker(s, 0.7, 0.6, "The system")
title(s, 1.05, "Two phases — all the power is offline")
def lane(y, label, sub, color, boxes):
    card(s, 0.7, y, 11.93, 1.72, PANEL)
    text(s, 1.0, y + 0.22, 3.0, 1.3, [{"runs": [(label, 16, color, True, False, HEAD)]},
                                      {"before": 4, "line": 1.05, "runs": [(sub, 11.5, SLATE)]}], anchor=MSO_ANCHOR.MIDDLE)
    bx, bw, bgap = 4.05, 2.5, 0.28
    for i, b in enumerate(boxes):
        x = bx + i * (bw + bgap)
        minibox(s, x, y + 0.42, bw, 0.88, b, color)
        if i < len(boxes) - 1:
            arrow(s, x + bw - 0.02, y + 0.42, bgap + 0.04, 0.88, color)
lane(2.15, "Pre-compute", "offline · LLM + GPU allowed", BRASS_D,
     ["Embed + features", "DeepSeek-V4\ngrades fit 0–5", "Kimi-K2.6\n2nd judge"])
lane(4.15, "Ranking step", "offline · CPU · no network · < 5 min", TEAL,
     ["Integrity Gate", "Score all\n100,000", "Top 100\n+ reasons"])
text(s, 0.7, 6.15, 11.93, 0.8,
     [{"runs": [("Spec §3 compliance:  ", 14, INK, True),
                ("no LLM, GPU or network during ranking — every model call lives in pre-compute and ships as a cached artifact the ranker reads.", 14, SLATE)]}])

# ============================================================ 8 — INTEGRITY GATE
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Layer 1 — the Integrity Gate")
title(s, 1.05, "Reject the impossible and the fake — first", 31)
det = [("Honeypot check", "‘Expert’ in a skill used 0 months · tenure longer than the whole career · contradictory dates. ~84 hits ≈ the ~80 planted honeypots."),
       ("Stuffer check", "Many AI skills, but the career narrative shows none of the work, and the title is wrong. Keywords without evidence don't count.")]
for i, (t, d) in enumerate(det):
    y = 2.15 + i * 1.5
    card(s, 0.7, y, 7.4, 1.32, PANEL)
    hallmark(s, 1.15, y + 0.66, 0.34, CLAY)
    text(s, 1.6, y + 0.24, 6.3, 0.5, [{"runs": [(t, 19, INK, True, False, HEAD)]}])
    text(s, 1.6, y + 0.72, 6.35, 0.55, [{"line": 1.08, "runs": [(d, 13.5, SLATE)]}])
card(s, 8.35, 2.15, 4.28, 2.82, INK)
text(s, 8.7, 2.45, 3.6, 0.5, [{"runs": [("VERDICT AT RUNTIME", 12, MIST, True)]}])
text(s, 8.7, 2.95, 3.6, 1.0, [{"runs": [("Flagged profiles are forced to the bottom — they can never reach the top 100.", 15, PAPER)]}])
text(s, 8.7, 3.72, 3.6, 1.1, [{"runs": [("0", 44, BRASS, True, False, HEAD)]},
                              {"before": 2, "runs": [("honeypots in our top 100", 13.5, MIST)]}])

# ============================================================ 9 — SEEING THROUGH THE DISGUISE
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Why it works — a concrete case")
title(s, 1.05, "Seeing through the disguise")
# fake
card(s, 0.7, 2.15, 5.83, 3.15, CLAY_LT)
text(s, 1.05, 2.4, 5.2, 0.4, [{"runs": [("THE DISGUISE", 12, CLAY, True)]}])
text(s, 1.05, 2.8, 5.2, 1.7,
     [{"line": 1.15, "after": 6, "runs": [("Marketing Manager", 17, INK, True, False, HEAD), (" · 9 yrs", 15, SLATE)]},
      {"line": 1.15, "after": 6, "runs": [("Skills: RAG, NLP, Embeddings, Ranking", 13.5, INK_SOFT)]},
      {"line": 1.15, "runs": [("Career: brand campaigns & strategy. No ML system ever built.", 13.5, SLATE)]}])
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.05, 4.6, 5.15, 0.5, fill=PAPER, line=CLAY, lw=1.25, radius=0.2)
text(s, 1.05, 4.6, 5.15, 0.5, [{"align": PP_ALIGN.CENTER, "runs": [("Judge verdict:  TIER 1  —  keywords, no evidence", 13, CLAY, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# gem
card(s, 6.8, 2.15, 5.83, 3.15, TEAL_LT)
text(s, 7.15, 2.4, 5.2, 0.4, [{"runs": [("THE REAL THING", 12, TEAL, True)]}])
text(s, 7.15, 2.8, 5.2, 1.7,
     [{"line": 1.15, "after": 6, "runs": [("Data Engineer", 17, INK, True, False, HEAD), (" · 7 yrs", 15, SLATE)]},
      {"line": 1.15, "after": 6, "runs": [("No ‘RAG’ or ‘Pinecone’ in the skills list.", 13.5, INK_SOFT)]},
      {"line": 1.15, "runs": [("Career: “built and shipped a product recommendation system at scale.”", 13.5, SLATE)]}])
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.15, 4.6, 5.15, 0.5, fill=PAPER, line=TEAL, lw=1.25, radius=0.2)
text(s, 7.15, 4.6, 5.15, 0.5, [{"align": PP_ALIGN.CENTER, "runs": [("Judge verdict:  TIER 5  —  real work, plainly told", 13, TEAL, True)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 5.65, 11.9, 0.5, [{"runs": [("Keyword counting ranks the fake #1 and never finds the gem. Parakh reads the work and flips both.", 15, INK, True, True)]}])

# ============================================================ 10 — WHAT THE TEACHER JUDGES
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Layer 2 — the distilled teacher")
title(s, 1.05, "What the teacher actually judges")
card(s, 0.7, 2.15, 5.83, 3.2, TEAL_LT)
text(s, 1.05, 2.45, 5.2, 0.5, [{"runs": [("Must-haves it rewards", 19, INK, True, False, HEAD)]}])
for i, t in enumerate(["Production retrieval / ranking / recsys, shipped to users",
                       "Embeddings + a vector DB (FAISS, Pinecone, Qdrant…)",
                       "Rigorous eval — NDCG, MRR, A/B testing",
                       "5–9 yrs at product companies; in / open to India"]):
    text(s, 1.35, 3.05 + i * 0.53, 5.0, 0.5, [{"runs": [("·  ", 15, TEAL, True), (t, 13.5, INK_SOFT)]}])
card(s, 6.8, 2.15, 5.83, 3.2, CLAY_LT)
text(s, 7.15, 2.45, 5.2, 0.5, [{"runs": [("Disqualifiers it catches", 19, INK, True, False, HEAD)]}])
for i, t in enumerate(["Stuffers — AI keywords with no supporting work",
                       "Pure research / academia, no production",
                       "Recent LangChain-only ‘AI experience’",
                       "Whole career at IT-services / consulting"]):
    text(s, 7.45, 3.05 + i * 0.53, 5.0, 0.5, [{"runs": [("·  ", 15, CLAY, True), (t, 13.5, INK_SOFT)]}])
text(s, 0.7, 5.65, 11.9, 0.5, [{"runs": [("Grades every plausible candidate 0–5 with a written reason — the rubric a senior recruiter carries in their head.", 14, SLATE)]}])

# ============================================================ 11 — QUORUM
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Layer 3 — agreement, not one opinion")
title(s, 1.05, "Two independent judges — the top is where they agree", 29)
card(s, 0.7, 2.35, 3.5, 1.5, BRASS_LT)
text(s, 1.0, 2.6, 2.9, 1.0, [{"runs": [("DeepSeek-V4", 18, INK, True, False, HEAD)]}, {"before": 3, "runs": [("primary teacher", 13, SLATE)]}])
card(s, 0.7, 4.05, 3.5, 1.5, TEAL_LT)
text(s, 1.0, 4.3, 2.9, 1.0, [{"runs": [("Kimi-K2.6", 18, INK, True, False, HEAD)]}, {"before": 3, "runs": [("independent 2nd judge", 13, SLATE)]}])
arrow(s, 4.35, 2.9, 0.9, 2.0, BRASS)
arrow(s, 4.35, 3.9, 0.9, 2.0, BRASS)
card(s, 5.35, 2.9, 3.3, 2.05, INK)
text(s, 5.65, 3.2, 2.7, 1.6, [{"runs": [("Average the tiers", 17, PAPER, True, False, HEAD)]},
                              {"before": 8, "line": 1.12, "runs": [("Only candidates ", 13.5, MIST), ("both", 13.5, BRASS, True), (" rate tier-5 rise to the very top.", 13.5, MIST)]}], anchor=MSO_ANCHOR.MIDDLE)
card(s, 9.0, 2.9, 3.63, 2.05, PANEL)
text(s, 9.3, 3.15, 3.1, 1.7, [{"runs": [("Why it pays off", 14, INK, True, False, HEAD)]},
                             {"before": 6, "line": 1.12, "runs": [("Different model families make different mistakes. Requiring agreement cancels each one's overconfidence.", 13, SLATE)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 5.85, 11.9, 0.5, [{"runs": [("An echo of my ", 13.5, SLATE), ("Quorum", 13.5, INK, True, True), (" project: trust rises where independent models converge.", 13.5, SLATE)]}])

# ============================================================ 12 — BEHAVIORAL
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Layer 4 — availability & red flags")
title(s, 1.05, "Beyond the résumé")
card(s, 0.7, 2.15, 5.83, 3.2, TEAL_LT)
text(s, 1.05, 2.45, 5.2, 0.5, [{"runs": [("Availability signals", 20, INK, True, False, HEAD)]}])
for i, t in enumerate(["Recruiter response rate — will they reply?",
                       "Last active — a 6-month ghost isn't hireable",
                       "Open-to-work + notice period",
                       "Interview completion, saved-by-recruiters"]):
    text(s, 1.35, 3.05 + i * 0.53, 5.0, 0.5, [{"runs": [("·  ", 15, TEAL, True), (t, 14, INK_SOFT)]}])
card(s, 6.8, 2.15, 5.83, 3.2, CLAY_LT)
text(s, 7.15, 2.45, 5.2, 0.5, [{"runs": [("Red flags it down-weights", 20, INK, True, False, HEAD)]}])
for i, t in enumerate(["Based abroad and won't relocate",
                       "Long notice period",
                       "Title-chaser — job-hops every ~18 months",
                       "‘Architect’ who hasn't coded in 18 months"]):
    text(s, 7.45, 3.05 + i * 0.53, 5.0, 0.5, [{"runs": [("·  ", 15, CLAY, True), (t, 14, INK_SOFT)]}])
text(s, 0.7, 5.65, 11.9, 0.5, [{"runs": [("A perfect-on-paper candidate who never replies isn't a hire. These signals sharpen the very top.", 14, SLATE)]}])

# ============================================================ 13 — RIGOR
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Rigor")
title(s, 1.05, "We didn't fly blind")
text(s, 0.7, 2.0, 11.9, 0.9,
     [{"line": 1.18, "runs": [("No public leaderboard and only three submissions — so most teams guess. We built our own eval harness and judged every version with a ", 15.5, SLATE),
                              ("third, independent model", 15.5, INK, True), (" that never touches the ranking.", 15.5, SLATE)]}])
judges = [("DeepSeek-V4", "teacher — grades 3,000", BRASS_D),
          ("Kimi-K2.6", "2nd judge — top 100", TEAL),
          ("Llama-3.3-70B", "independent scorer", INK_SOFT)]
for i, (m, role, c) in enumerate(judges):
    x = 0.7 + i * 4.05
    card(s, x, 3.25, 3.78, 1.55, PANEL)
    hallmark(s, x + 0.55, 3.62, 0.32, c)
    text(s, x + 1.0, 3.5, 2.6, 0.5, [{"runs": [(m, 17, c, True, False, HEAD)]}])
    text(s, x + 0.35, 4.05, 3.2, 0.6, [{"runs": [(role, 13.5, SLATE)]}])
text(s, 0.7, 5.2, 11.9, 0.6, [{"runs": [("Rule: ", 14, INK, True), ("a change shipped only if the independent judge scored it higher. Nothing made the cut on faith.", 14, SLATE)]}])

# ============================================================ 14 — RESULTS
s = slide(PAPER)
kicker(s, 0.7, 0.6, "Results")
title(s, 1.05, "Every layer earned its place", 34)
text(s, 0.7, 1.95, 8.0, 0.5, [{"runs": [("NDCG@10, scored by an independent third model (Llama-3.3-70B)", 14.5, SLATE)]}])
base_y, max_h = 5.75, 3.05
bars = [("Rule-only", 0.64, MIST), ("+ Teacher", 0.88, BRASS), ("+ Ensemble", 1.00, TEAL)]
bx0, bw, bgap = 1.2, 1.8, 1.05
for i, (lab, val, col) in enumerate(bars):
    x = bx0 + i * (bw + bgap); h = max_h * val
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, base_y - h, bw, h, fill=col, radius=0.04)
    text(s, x - 0.2, base_y - h - 0.55, bw + 0.4, 0.5, [{"align": PP_ALIGN.CENTER, "runs": [(f"{val:.2f}", 26, INK, True, False, HEAD)]}])
    text(s, x - 0.2, base_y + 0.1, bw + 0.4, 0.4, [{"align": PP_ALIGN.CENTER, "runs": [(lab, 14, INK_SOFT, True)]}])
bl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(base_y), Inches(8.4), Inches(0.018))
bl.fill.solid(); bl.fill.fore_color.rgb = SLATE; bl.line.fill.background(); bl.shadow.inherit = False
card(s, 9.7, 2.35, 2.95, 1.35, INK)
text(s, 9.95, 2.55, 2.5, 1.0, [{"runs": [("10 / 10", 30, BRASS, True, False, HEAD)]}, {"runs": [("top-10 are tier-5 by both judges", 12.5, MIST)]}])
card(s, 9.7, 3.85, 2.95, 1.35, PANEL)
text(s, 9.95, 4.05, 2.5, 1.0, [{"runs": [("0.98", 30, TEAL, True, False, HEAD)]}, {"runs": [("NDCG@50 (ensemble)", 12.5, SLATE)]}])
text(s, 0.7, 6.55, 11.9, 0.6, [{"runs": [("Honest read: ", 12.5, INK, True), ("three different models agree the top-10 is excellent — a strong proxy, not the organisers' hidden ground truth.", 12.5, SLATE, False, True)]}])

# ============================================================ 15 — WHY EVOLVE WINS
s = slide(PAPER)
kicker(s, 0.7, 0.6, "The bottom line")
title(s, 1.05, "Why this is different")
card(s, 0.7, 2.1, 5.83, 3.5, PANEL)
text(s, 1.05, 2.35, 5.2, 0.5, [{"runs": [("Everyone else", 18, SLATE, True, False, HEAD)]}])
for i, t in enumerate(["Embed everything, cosine, sort", "Fooled by keyword-stuffers", "Buries plain-language gems",
                       "Risks honeypot disqualification", "Flies blind — no way to measure"]):
    text(s, 1.05, 2.95 + i * 0.5, 5.2, 0.45, [{"runs": [("✕  ", 13, CLAY, True), (t, 14, INK_SOFT)]}])
card(s, 6.8, 2.1, 5.83, 3.5, TEAL_LT)
text(s, 7.15, 2.35, 5.2, 0.5, [{"runs": [("Parakh / Team Evolve", 18, INK, True, False, HEAD)]}])
for i, t in enumerate(["Assays every profile before trusting it", "Distills a frontier LLM into a legal offline ranker",
                       "Two independent judges must agree at the top", "0 honeypots · top-10 all tier-5",
                       "Proven by a third, independent model"]):
    text(s, 7.15, 2.95 + i * 0.5, 5.35, 0.45, [{"runs": [("✓  ", 13, TEAL, True), (t, 14, INK)]}])
text(s, 0.7, 5.85, 11.9, 0.5, [{"runs": [("Same rules as everyone. A fundamentally different bet: understand the candidate, don't pattern-match the keywords.", 14, INK, True, True)]}])

# ============================================================ 16 — CLOSE
s = slide(INK)
for (cx, cy, d) in [(1.7, 6.0, 2.6), (0.8, 1.4, 1.6)]:
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    ring.fill.background(); ring.line.color.rgb = INK_SOFT; ring.line.width = Pt(1.5); ring.shadow.inherit = False
hallmark(s, 6.67, 2.2, 0.5, BRASS)
text(s, 0, 2.6, SW, 1.0, [{"align": PP_ALIGN.CENTER, "runs": [("Parakh", 58, PAPER, True, False, HEAD)]}])
text(s, 0, 3.75, SW, 0.6, [{"align": PP_ALIGN.CENTER, "runs": [("Assaying talent — genuine fits, surfaced and trusted.", 19, BRASS, False, True, HEAD)]}])
bx = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.97, 4.7, 5.4, 0.72, fill=INK_SOFT, radius=0.2)
text(s, 3.97, 4.7, 5.4, 0.72, [{"align": PP_ALIGN.CENTER, "runs": [("python rank.py --candidates candidates.jsonl", 14, MIST, False, False, MONO)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0, 6.05, SW, 0.5, [{"align": PP_ALIGN.CENTER, "runs": [("Team Evolve", 15, PAPER, True)]}])
text(s, 0, 6.5, SW, 0.5, [{"align": PP_ALIGN.CENTER, "runs": [(CREDITS, 13, MIST)]}])

prs.save("parakh_deck.pptx")
print("saved parakh_deck.pptx —", len(prs.slides._sldIdLst), "slides")
