#!/usr/bin/env python3
"""Insert a 'Why not just embeddings?' rationale slide into the official deck,
matching the template's look (background reused from an existing content slide)."""
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK = RGBColor(0x20, 0x21, 0x24)
PUR = RGBColor(0x6D, 0x28, 0xD9)
TITLE_GREY = RGBColor(0x3C, 0x40, 0x43)
F = "Arial"

prs = Presentation("Parakh_official.pptx")
S = list(prs.slides)

# background image from an existing content slide (slide 4 = Ranking Methodology)
bg_blob = None
for sh in S[3].shapes:
    if sh.shape_type == 13:  # PICTURE
        bg_blob = sh.image.blob
        break
assert bg_blob, "no background picture found"

blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
s = prs.slides.add_slide(blank)
s.shapes.add_picture(io.BytesIO(bg_blob), 0, 0, width=prs.slide_width, height=prs.slide_height)

def text(x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(para.get("after", 0))
        p.line_spacing = para.get("line", 1.03)
        for t, sz, c, b in para["runs"]:
            r = p.add_run(); r.text = t
            r.font.name = F; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b

# title in template style
text(0.53, 0.89, 8.9, 0.51, [{"runs": [("Design Rationale - Why Not Just Embeddings?", 18, TITLE_GREY, False)]}])

qa = [
    ("Where would the labels come from?",
     "A learned ranker needs training labels - and this dataset ships none. The only way to get a supervision signal is to synthesize judgment, which is exactly what our offline teachers produce. Any 'embeddings + learned ranker' build starts by bootstrapping labels the same way we did; we simply spent the budget on making that judgment reliable (three model families) instead of hiding it inside a trained model."),
    ("The rules reward this design",
     "Runtime is CPU-only, no network, five minutes. The challenge brief itself says: 'plan for a small ranker over precomputed features.' Parakh is that sentence taken seriously - think offline with frontier models once, serve cached judgment in milliseconds, reproduce byte-for-byte."),
    ("On this dataset, cosine similarity rewards the fraud",
     "Keyword-stuffed profiles sit close to the JD in embedding space precisely because they copied its vocabulary - semantic similarity scores the disguise, not the work. Verified narrative evidence plus judge agreement does not. Embedding retrieval is on our roadmap for 200k+ scale - as a recall layer underneath the same Integrity Gate and judges, never as a replacement for them."),
]
first = True
tb = s.shapes.add_textbox(Inches(0.41), Inches(1.55), Inches(9.32), Inches(3.7))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
for q, a in qa:
    pq = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    pq.space_after = Pt(1)
    r = pq.add_run(); r.text = q
    r.font.name = F; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = PUR
    pa = tf.add_paragraph(); pa.space_after = Pt(7); pa.line_spacing = 1.03
    r = pa.add_run(); r.text = a
    r.font.name = F; r.font.size = Pt(10.5); r.font.color.rgb = INK

# move the new slide (last) to position 5 (0-based index 4), right after Ranking Methodology
lst = prs.slides._sldIdLst
ids = list(lst)
lst.remove(ids[-1])
lst.insert(4, ids[-1])

prs.save("Parakh_official.pptx")
print("inserted rationale slide at position 5; total", len(prs.slides._sldIdLst))
